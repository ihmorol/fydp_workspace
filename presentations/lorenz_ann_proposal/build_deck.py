from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
OUTPUT_PATH = OUT_DIR / "lorenz_ann_supervisor_informal.pptx"
UIU_LOGO = OUT_DIR.parents[1] / "fydp_template" / "uiu.png"

NAVY = RGBColor(16, 35, 60)
GOLD = RGBColor(212, 167, 44)
CREAM = RGBColor(246, 243, 235)
INK = RGBColor(33, 43, 54)
SLATE = RGBColor(82, 96, 113)
MIST = RGBColor(226, 231, 236)
TEAL = RGBColor(22, 108, 122)
RUST = RGBColor(184, 88, 65)
SAGE = RGBColor(110, 132, 102)
WHITE = RGBColor(255, 255, 255)

DECK_LABEL = "Informal Supervisor Deck"


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size=20,
    bold=False,
    color=INK,
    font_name="Arial",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.08,
    italic=False,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign

    lines = str(text).split("\n")
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line if line else " "
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        run = p.runs[0]
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, *, font_size=21, color=INK, level_indent=0.23):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        p.space_before = Pt(0)
        p.left_margin = Inches(level_indent)
        p.hanging = Inches(0.12)
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, *, fill, line=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1.2)
    return shape


def add_card(slide, left, top, width, height, title, body, *, accent, body_size=18):
    add_rect(slide, left, top, width, height, fill=WHITE, line=MIST, radius=True)
    add_rect(slide, left, top, Inches(0.12), height, fill=accent, line=accent)
    add_textbox(slide, left + Inches(0.24), top + Inches(0.16), width - Inches(0.38), Inches(0.36), title, font_size=20, bold=True, color=INK)
    add_textbox(slide, left + Inches(0.24), top + Inches(0.56), width - Inches(0.38), height - Inches(0.72), body, font_size=body_size, color=SLATE)


def set_background(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_header(slide, section_no, title):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.62), fill=NAVY, line=NAVY)
    add_rect(slide, Inches(0), Inches(0.62), Inches(13.333), Inches(0.05), fill=GOLD, line=GOLD)
    add_rect(slide, Inches(0.42), Inches(0.16), Inches(0.62), Inches(0.32), fill=GOLD, line=GOLD, radius=True)
    add_textbox(slide, Inches(0.42), Inches(0.16), Inches(0.62), Inches(0.32), section_no, font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(1.15), Inches(0.11), Inches(8.5), Inches(0.44), title, font_size=26, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(9.55), Inches(0.12), Inches(3.3), Inches(0.34), DECK_LABEL, font_size=12, color=WHITE, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def add_footer(slide, text, *, dark=False):
    color = WHITE if dark else SLATE
    add_textbox(slide, Inches(0.45), Inches(7.02), Inches(12.1), Inches(0.22), text, font_size=10, color=color)


def add_logo(slide, *, dark=False):
    if UIU_LOGO.exists():
        width = Inches(0.78) if not dark else Inches(1.0)
        slide.shapes.add_picture(str(UIU_LOGO), Inches(11.82), Inches(6.5), width=width)


def add_divider_line(slide, left, top, width, *, color=GOLD, line_width=2.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(line_width)
    return line


def title_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)

    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=NAVY, line=NAVY)
    add_rect(slide, Inches(-0.2), Inches(5.92), Inches(8.25), Inches(1.88), fill=GOLD, line=GOLD)
    add_textbox(slide, Inches(0.72), Inches(0.84), Inches(9.0), Inches(1.55), "Solving the Lorenz-1960 ODE System\nUsing ANN and Identifying the Optimal Architecture", font_size=28, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.78), Inches(2.62), Inches(8.0), Inches(0.68), "Informal research proposal for supervisor discussion", font_size=18, color=CREAM)
    add_divider_line(slide, Inches(0.78), Inches(3.42), Inches(4.3), color=GOLD, line_width=2.5)
    add_textbox(slide, Inches(0.78), Inches(3.62), Inches(6.5), Inches(0.3), "Department of CSE, United International University", font_size=15, color=CREAM)
    add_textbox(slide, Inches(0.78), Inches(3.92), Inches(5.4), Inches(0.28), "FYDP research proposal meeting", font_size=15, color=CREAM)

    add_rect(slide, Inches(0.82), Inches(5.15), Inches(8.7), Inches(1.18), fill=WHITE, line=WHITE, radius=True)
    add_textbox(slide, Inches(1.05), Inches(5.42), Inches(8.15), Inches(0.64), "Problem statement: Can a plain ANN learn the Lorenz-1960 solution accurately, and which ANN architecture works best when compared against reliable numerical and PINN-style references?", font_size=20, bold=True, color=NAVY)

    if UIU_LOGO.exists():
        slide.shapes.add_picture(str(UIU_LOGO), Inches(11.25), Inches(5.82), width=Inches(1.18))
    add_textbox(slide, Inches(10.15), Inches(6.97), Inches(2.45), Inches(0.2), "March 13, 2026", font_size=10, color=CREAM, align=PP_ALIGN.RIGHT)
    return slide


def background_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, CREAM)
    add_header(slide, "01", "Background")

    add_textbox(slide, Inches(0.68), Inches(0.96), Inches(11.9), Inches(0.44), "Ordinary differential equations describe the evolution of many physical systems, and the Lorenz-1960 model is a compact nonlinear benchmark that is already used in the literature.", font_size=20)

    add_card(slide, Inches(0.82), Inches(1.8), Inches(3.85), Inches(3.9), "ODE Context", "ODEs appear in weather dynamics, fluid interaction, control systems, and many other time-evolving physical processes.", accent=NAVY, body_size=18)
    add_card(slide, Inches(4.86), Inches(1.8), Inches(3.85), Inches(3.9), "Lorenz-1960 Benchmark", "This project uses the three interdependent Lorenz-1960 equations from Paper 1 Section 4.2 with k=2, l=1 and fixed initial conditions.", accent=TEAL, body_size=18)
    add_card(slide, Inches(8.9), Inches(1.8), Inches(3.6), Inches(3.9), "Reference Literature", "Paper 1 provides the benchmark problem, Paper 2 positions ANN-based differential equation solvers, and Paper 3 gives a contrasting Neural ODE viewpoint.", accent=RUST, body_size=18)

    add_footer(slide, "Background slide keeps the research grounded in the benchmark problem and literature context.")
    add_logo(slide)
    return slide


def motivation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, CREAM)
    add_header(slide, "02", "Motivation")

    add_card(slide, Inches(0.82), Inches(1.08), Inches(3.8), Inches(4.95), "Why Numerical Solvers Are Not the Full Story", "Classical methods such as Runge-Kutta are accurate, but they solve the system again for each experiment and do not directly produce a learned surrogate.", accent=NAVY, body_size=19)
    add_card(slide, Inches(4.78), Inches(1.08), Inches(3.8), Inches(4.95), "Why ANN Is Interesting", "A trained ANN can approximate the solution map quickly after training, which is attractive when repeated evaluation or future extension to multiple settings is needed.", accent=TEAL, body_size=19)
    add_card(slide, Inches(8.74), Inches(1.08), Inches(3.78), Inches(4.95), "Why This Benchmark", "Lorenz-1960 is small enough to study carefully but nonlinear enough to make architecture choice and solver comparison meaningful.", accent=RUST, body_size=19)

    add_rect(slide, Inches(0.92), Inches(6.2), Inches(11.45), Inches(0.52), fill=WHITE, line=MIST, radius=True)
    add_textbox(slide, Inches(1.1), Inches(6.32), Inches(11.0), Inches(0.24), "Motivation in one line: build a reliable numerical baseline now, then use it to judge which ANN architecture is actually worth trusting later.", font_size=19, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_footer(slide, "The motivation is practical: accuracy, reusability, and future comparison on a fixed benchmark.")
    add_logo(slide)
    return slide


def contribution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, CREAM)
    add_header(slide, "03", "Expected Contribution")

    contributions = [
        ("1", "Problem-Focused Benchmark", "A clean FYDP benchmark centered on the Lorenz-1960 ODE system rather than a generic ANN demo.", NAVY),
        ("2", "Architecture Study Direction", "A research path aimed at discovering how depth, width, activation, and training strategy affect solution quality.", TEAL),
        ("3", "Comparison Framework", "A common evaluation frame for numerical solvers, plain ANN, and PINN-style references on the same problem.", RUST),
        ("4", "Reusable Workflow", "A reproducible pipeline for literature review, solver implementation, validation, and later ANN experiments.", SAGE),
    ]

    positions = [(0.85, 1.6), (6.75, 1.6), (0.85, 4.12), (6.75, 4.12)]
    for (left, top), (code, title, body, accent) in zip(positions, contributions):
        add_rect(slide, Inches(left), Inches(top), Inches(5.7), Inches(1.95), fill=WHITE, line=MIST, radius=True)
        add_rect(slide, Inches(left + 0.2), Inches(top + 0.22), Inches(0.58), Inches(0.58), fill=accent, line=accent, radius=True)
        add_textbox(slide, Inches(left + 0.2), Inches(top + 0.22), Inches(0.58), Inches(0.58), code, font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(left + 0.95), Inches(top + 0.18), Inches(4.4), Inches(0.34), title, font_size=21, bold=True, color=INK)
        add_textbox(slide, Inches(left + 0.95), Inches(top + 0.62), Inches(4.45), Inches(0.95), body, font_size=18, color=SLATE)

    add_footer(slide, "Contribution is framed as a benchmark and methodology contribution, not an overclaimed final result.")
    add_logo(slide)
    return slide


def novelty_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, CREAM)
    add_header(slide, "04", "How Is This a Novel Problem?")

    add_textbox(slide, Inches(0.72), Inches(0.98), Inches(11.7), Inches(0.42), "The novelty is not the existence of ANN solvers alone. The novelty is the specific benchmark, comparison frame, and architecture question being asked together.", font_size=20)

    add_card(slide, Inches(0.84), Inches(1.8), Inches(3.85), Inches(4.4), "What Existing Work Already Does", "Paper 1 solves Lorenz-1960 with DeepONet-style operator learning. Paper 2 reviews ANN-based differential equation methods. Paper 3 introduces Neural ODEs as a different paradigm.", accent=SLATE, body_size=17)
    add_card(slide, Inches(4.9), Inches(1.8), Inches(3.85), Inches(4.4), "What Is Missing", "The literature does not clearly answer which plain ANN architecture is best for solving the Lorenz-1960 benchmark under a unified experimental setup.", accent=TEAL, body_size=18)
    add_card(slide, Inches(8.96), Inches(1.8), Inches(3.45), Inches(4.4), "Why That Gap Matters", "Without this architecture-level evidence, it is hard to justify when a plain ANN is sufficient and when a more structured neural solver is necessary.", accent=RUST, body_size=18)

    add_footer(slide, "Novelty comes from the research question and benchmark framing, not from claiming ANN-for-ODE is entirely new.")
    add_logo(slide)
    return slide


def complex_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, CREAM)
    add_header(slide, "05", "How Is This a Complex Engineering Problem?")

    add_card(slide, Inches(0.82), Inches(1.02), Inches(5.75), Inches(5.55), "Why It Is Complex", "The problem combines mathematical modeling, numerical accuracy, software implementation, model evaluation, and design trade-off decisions. There is no single closed-form shortcut that solves all of these at once.", accent=NAVY, body_size=20)

    add_bullets(slide, Inches(6.88), Inches(1.35), Inches(5.3), Inches(4.8), [
        "The Lorenz-1960 system is nonlinear and coupled, so small implementation mistakes can propagate across all three states.",
        "A trustworthy ANN study needs a validated numerical reference first, otherwise model claims are weak.",
        "Architecture choice creates a multi-factor design problem involving depth, width, activation, optimizer, and stability.",
        "The final answer must balance accuracy, computational cost, reproducibility, and clarity of comparison.",
    ], font_size=19)

    add_rect(slide, Inches(6.92), Inches(5.86), Inches(5.35), Inches(0.52), fill=WHITE, line=MIST, radius=True)
    add_textbox(slide, Inches(7.12), Inches(5.99), Inches(4.95), Inches(0.24), "Engineering complexity = model design + solver validation + evidence-based comparison", font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_footer(slide, "This slide is tuned for supervisor review: it explains why the task qualifies as an engineering problem, not just a coding exercise.")
    add_logo(slide)
    return slide


def fydp_scope_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, CREAM)
    add_header(slide, "06", "FYDP-I Scope")

    add_card(slide, Inches(0.82), Inches(1.0), Inches(5.55), Inches(5.6), "What FYDP-I Will Cover", "Complete the literature review.\nImplement the Lorenz-1960 equations using a Python ODE solver library.\nImplement a custom Runge-Kutta baseline.\nOptimize and validate the numerical reference outputs for later ANN comparison.", accent=NAVY, body_size=20)
    add_card(slide, Inches(6.65), Inches(1.0), Inches(5.7), Inches(2.5), "Immediate Deliverables", "Literature review matrix and gap analysis\nVerified numerical solution plots\nReusable solver code and baseline evaluation notes", accent=TEAL, body_size=19)
    add_card(slide, Inches(6.65), Inches(3.85), Inches(5.7), Inches(2.75), "What Is Not Promised Yet", "Final ANN architecture optimization\nFull ANN versus PINN experiment set\nFinal recommendation across all solver families", accent=RUST, body_size=19)

    add_footer(slide, "FYDP-I is positioned as a strong foundation phase: literature plus validated numerical benchmark.")
    add_logo(slide)
    return slide


def closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)

    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill=NAVY, line=NAVY)
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.78), Inches(12.05), Inches(5.9))
    frame.fill.background()
    frame.line.color.rgb = GOLD
    frame.line.width = Pt(3)

    add_textbox(slide, Inches(1.15), Inches(1.3), Inches(8.4), Inches(0.58), "Discussion / Next Step", font_size=22, bold=True, color=GOLD)
    add_textbox(slide, Inches(1.15), Inches(1.95), Inches(10.4), Inches(1.7), "The proposal is simple to state but strong enough to grow: first validate the Lorenz-1960 numerical benchmark, then use that benchmark to study and compare ANN architectures in later phases.", font_size=28, bold=True, color=WHITE)
    add_textbox(slide, Inches(1.18), Inches(4.2), Inches(9.45), Inches(0.82), "Supervisor feedback needed on FYDP-I boundary, preferred baseline solver, and how broad the later ANN search should be.", font_size=20, color=CREAM)
    add_textbox(slide, Inches(1.18), Inches(5.72), Inches(3.8), Inches(0.28), "Suggested discussion points", font_size=16, bold=True, color=GOLD)
    add_textbox(slide, Inches(1.18), Inches(6.0), Inches(8.9), Inches(0.34), "scope approval, numerical baseline standard, and whether a small ANN prototype is needed in FYDP-I", font_size=16, color=CREAM)
    if UIU_LOGO.exists():
        slide.shapes.add_picture(str(UIU_LOGO), Inches(11.2), Inches(5.3), width=Inches(1.22))
    add_footer(slide, "Prepared from current FYDP notes, literature summary, and supervisor-facing scope clarification.", dark=True)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.title = "Informal Supervisor Research Proposal Deck"
    prs.core_properties.subject = "Lorenz-1960 ANN proposal"

    title_problem_slide(prs)
    background_slide(prs)
    motivation_slide(prs)
    contribution_slide(prs)
    novelty_slide(prs)
    complex_problem_slide(prs)
    fydp_scope_slide(prs)
    closing_slide(prs)

    prs.save(OUTPUT_PATH)
    print(f"Created {OUTPUT_PATH}")


if __name__ == "__main__":
    build()
